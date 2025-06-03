from crewai.tools import BaseTool
from typing import Type, List, Union
from pydantic import BaseModel, Field
from pypdf import PdfReader
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
import json
import os

class ParseDocumentInput(BaseModel):
    """Input schema for document parsing tool."""
    file_paths: Union[str, List[str]] = Field(..., description="Path or list of paths to document files to parse")

class WebResearchInput(BaseModel):
    """Input schema for web research tool."""
    query: str = Field(..., description="Search query about the startup or industry")

class DocumentParserTool(BaseTool):
    name: str = "Document Parser"
    description: str = (
        "Tool for parsing pitch deck documents (PDF/PPT) and additional files"
    )
    args_schema: Type[BaseModel] = ParseDocumentInput

    def _run(self, file_paths: Union[str, List[str]]) -> str:
        # Convert single file path to list for consistent handling
        if isinstance(file_paths, str):
            file_paths = [file_paths]
        
        all_text = []
        failed_files = []
        
        for file_path in file_paths:
            try:
                print(f"\nProcessing file: {file_path}")
                
                # Convert to absolute path if not already
                abs_file_path = os.path.abspath(file_path)
                print(f"Absolute path: {abs_file_path}")
                
                # Verify file exists
                exists = os.path.exists(abs_file_path)
                print(f"File exists: {exists}")
                if not exists:
                    raise FileNotFoundError(f"File not found: {abs_file_path}")
                
                # Check file extension
                file_ext = os.path.splitext(abs_file_path.lower())[1]
                if file_ext == '.pdf':
                    text = self._parse_pdf(abs_file_path)
                elif file_ext in ['.ppt', '.pptx']:
                    text = self._parse_ppt(abs_file_path)
                else:
                    raise ValueError(f"Unsupported file format: {file_ext}. Only PDF and PPT/PPTX are supported.")
                
                # Add file header and content
                file_name = os.path.basename(file_path)
                all_text.append(f"\n{'='*20} Document: {file_name} {'='*20}\n{text}")
                print(f"Successfully processed: {file_name}")
            
            except Exception as e:
                error_msg = f"Error processing {os.path.basename(file_path)}: {str(e)}"
                print(f"WARNING: {error_msg}")
                failed_files.append(error_msg)
                continue
        
        if not all_text:
            raise ValueError("No documents were successfully processed. " + 
                           "Failed files:\n" + "\n".join(failed_files))
        
        # Combine all document texts with processing summary
        result = []
        if failed_files:
            result.append("WARNING: Some files failed to process:\n" + 
                        "\n".join(f"- {err}" for err in failed_files))
        
        result.extend(all_text)
        return "\n\n".join(result)

    def _parse_pdf(self, file_path: str) -> str:
        try:
            reader = PdfReader(file_path)
            text = []
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text.strip():  # Only include non-empty pages
                    text.append(f"Page {i+1}:\n{page_text}")
            return "\n\n".join(text)
        except Exception as e:
            raise ValueError(f"Error parsing PDF: {str(e)}")

    def _parse_ppt(self, file_path: str) -> str:
        try:
            prs = Presentation(file_path)
            text = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:  # Only include slides with content
                    text.append(f"Slide {i+1}:\n{' '.join(slide_text)}")
            return "\n\n".join(text)
        except Exception as e:
            raise ValueError(f"Error parsing PPT/PPTX: {str(e)}")

class WebResearchTool(BaseTool):
    name: str = "Web Research"
    description: str = (
        "Tool for conducting web research about startups and industries"
    )
    args_schema: Type[BaseModel] = WebResearchInput

    def _run(self, query: str) -> str:
        try:
            # This is a simple example - in production you might want to use 
            # more sophisticated search APIs and handle rate limiting
            search_url = f"https://api.bing.microsoft.com/v7.0/search"
            headers = {
                "Ocp-Apim-Subscription-Key": "YOUR-KEY-HERE"  # Replace with actual key
            }
            params = {
                "q": query,
                "count": 5,
                "offset": 0,
                "mkt": "en-US"
            }
            
            response = requests.get(search_url, headers=headers, params=params)
            results = response.json()
            
            if "webPages" in results:
                pages = results["webPages"]["value"]
                research_results = []
                
                for page in pages:
                    try:
                        # Get the webpage content
                        page_response = requests.get(page["url"])
                        soup = BeautifulSoup(page_response.text, 'html.parser')
                        
                        # Extract main content (this is a simple example)
                        content = " ".join([p.get_text() for p in soup.find_all('p')])
                        
                        research_results.append({
                            "title": page["name"],
                            "url": page["url"],
                            "snippet": page["snippet"],
                            "content": content[:500]  # First 500 chars of content
                        })
                    except Exception as e:
                        continue
                        
                return json.dumps(research_results, indent=2)
            
            return "No results found"
            
        except Exception as e:
            return f"Error performing web research: {str(e)}"
